# SPDX-License-Identifier: MIT
# Copyright (C) 2026 ThreadWeave contributors
"""Tests for xlsx/pptx text extraction in the SharePoint processor."""

import io
import sys

import pytest

sys.path.insert(0, "src")
from threadweave.connectors.sharepoint.processor import (
    DocumentProcessor,
    OPENPYXL_AVAILABLE,
    PPTX_AVAILABLE,
)

pytestmark = pytest.mark.skipif(
    not (OPENPYXL_AVAILABLE and PPTX_AVAILABLE),
    reason="openpyxl / python-pptx not installed",
)


def make_xlsx() -> bytes:
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Decisions"
    ws.append(["Topic", "Decision"])
    ws.append(["Session cache", "We decided to use Redis for the session cache"])
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def make_pptx() -> bytes:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Infrastructure decision"
    slide.placeholders[1].text = "We decided to standardize on PostgreSQL"
    slide.notes_slide.notes_text_frame.text = "Confirmed with engineering"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def get_processor() -> DocumentProcessor:
    proc = DocumentProcessor.__new__(DocumentProcessor)
    proc.temp_dir = "."
    return proc


def test_extract_xlsx_returns_cell_text_with_sheet_header():
    text = get_processor()._extract_xlsx(make_xlsx())
    assert "[Sheet: Decisions]" in text
    assert "We decided to use Redis for the session cache" in text


def test_extract_pptx_returns_slide_and_notes_text():
    text = get_processor()._extract_pptx(make_pptx())
    assert "[Slide 1]" in text
    assert "We decided to standardize on PostgreSQL" in text
    assert "[notes] Confirmed with engineering" in text


def test_extract_text_dispatches_xlsx_and_pptx():
    proc = get_processor()
    x = proc._extract_text(make_xlsx(), ".xlsx", "decisions.xlsx")
    assert "Redis" in x
    p = proc._extract_text(make_pptx(), ".pptx", "deck.pptx")
    assert "PostgreSQL" in p


def test_unsupported_extension_returns_empty():
    assert get_processor()._extract_text(b"\x00\x01", ".one", "notes.one") == ""
