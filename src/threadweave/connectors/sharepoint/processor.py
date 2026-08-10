# SPDX-License-Identifier: MIT
# Copyright (C) 2026 ThreadWeave contributors
"""
SharePoint Document Processor — mines documents into MemPalace.

Handles:
    - Downloading files from SharePoint via Graph API
    - Text extraction (native Office formats via Graph, PDFs via PyMuPDF)
    - Mapping SharePoint structure -> MemPalace (site -> wing, library -> room)
    - Batch and incremental processing
    - Processing status tracking

The processor uses Microsoft Graph to download documents and
MemPalace's mine pipeline for text ingestion. For binary formats,
it can optionally use the Graph API's built-in text extraction
or fall back to local extraction via PyMuPDF / python-docx.
"""

from __future__ import annotations

import asyncio
import hashlib
import io
import json
import logging
import os
import tempfile
from dataclasses import dataclass, field
from datetime import datetime, timezone
from pathlib import Path
from typing import Optional

MAX_FOLDER_DEPTH = 8  # recursion bound for full-drive imports

from threadweave.connectors.sharepoint.watcher import (
    GraphClient,
    ChangeNotification,
    SiteInfo,
    DriveInfo,
)

logger = logging.getLogger(__name__)

# Optional dependencies for binary document extraction
try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False

try:
    import docx  # python-docx
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import openpyxl
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

try:
    import pptx  # python-pptx
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


# ---- Data Classes ----

@dataclass
class ProcessedDocument:
    """Result of processing a single SharePoint document."""
    file_name: str
    file_path: str       # SharePoint path
    site_name: str
    library_name: str
    mime_type: str
    size_bytes: int
    text_content: str    # Extracted text
    word_count: int
    drawer_ids: list[str] = field(default_factory=list)  # MemPalace drawer IDs
    status: str = "processed"  # processed, skipped, error
    error: str = ""


@dataclass
class ProcessingBatch:
    """Result of processing a batch of documents."""
    site_id: str
    site_name: str
    drive_id: str
    drive_name: str
    documents: list[ProcessedDocument] = field(default_factory=list)
    started_at: str = ""
    completed_at: str = ""
    total_processed: int = 0
    total_errors: int = 0
    total_skipped: int = 0


# ---- Document Processor ----

class DocumentProcessor:
    """
    Downloads SharePoint documents and mines them into MemPalace.

    Mapping:
        SharePoint site    -> MemPalace wing
        Document library   -> MemPalace room
        Folder hierarchy   -> Room keywords / sub-rooms

    File type support:
        - Text/CSV/JSON   -> Direct ingestion
        - .docx            -> python-docx extraction
        - .pdf             -> PyMuPDF extraction
        - .xlsx            -> Basic via Graph API preview
        - .pptx            -> Basic via Graph API preview
        - Other            -> Skipped (configurable)
    """

    SUPPORTED_EXTENSIONS = {
        ".txt", ".csv", ".json", ".md", ".yaml", ".yml",
        ".py", ".js", ".ts", ".html", ".css", ".xml",
        ".docx", ".pdf", ".xlsx", ".pptx",
        ".odt", ".ods", ".odp",
        ".vsdx",
        ".mp4", ".mkv", ".mov", ".webm", ".mp3", ".wav", ".m4a",
        ".log", ".cfg", ".ini", ".toml",
    }

    def __init__(
        self,
        graph_client: GraphClient,
        mempalace_palace_path: str = "~/.mempalace/palace",
        temp_dir: str | None = None,
        max_file_size_mb: int = 50,
    ):
        self.graph = graph_client
        self.palace_path = os.path.expanduser(mempalace_palace_path)
        self.temp_dir = temp_dir or os.path.join(
            tempfile.gettempdir(), "threadweave_sharepoint"
        )
        self.max_file_size = max_file_size_mb * 1024 * 1024
        self._processed_hashes: set[str] = set()  # De-duplication
        self.stats = {
            "total_documents": 0,
            "total_processed": 0,
            "total_errors": 0,
            "total_skipped": 0,
            "total_bytes_ingested": 0,
        }

        os.makedirs(self.temp_dir, exist_ok=True)

    # ---- Main Processing Pipeline ----

    async def process_notification(
        self, notification: ChangeNotification
    ) -> ProcessedDocument | None:
        """
        Process a single change notification — download and mine.

        Returns:
            ProcessedDocument on success, None if skipped.
        """
        try:
            # Download the file
            content = await self.graph.download_file(
                site_id=notification.site_id,
                drive_id=notification.list_id,
                item_id=notification.item_id,
            )

            if len(content) > self.max_file_size:
                logger.info(
                    "Skipping large file: %s bytes", len(content)
                )
                self.stats["total_skipped"] += 1
                return None

            # Get file metadata
            file_name = self._extract_filename(notification)

            # Compute content hash for de-duplication
            content_hash = hashlib.sha256(content).hexdigest()
            if content_hash in self._processed_hashes:
                logger.debug("Skipping duplicate: %s", file_name)
                self.stats["total_skipped"] += 1
                return None
            self._processed_hashes.add(content_hash)

            # Extract text
            ext = Path(file_name).suffix.lower()
            text = self._extract_text(content, ext, file_name)

            if not text.strip():
                logger.info("No extractable text in: %s", file_name)
                self.stats["total_skipped"] += 1
                return None

            # Mine into MemPalace
            drawer_ids = await self._mine_to_mempalace(
                text=text,
                wing=self._sanitize_wing(notification.site_id),
                room=self._sanitize_room(notification.list_id),
                source_file=file_name,
            )

            self.stats["total_processed"] += 1
            self.stats["total_bytes_ingested"] += len(content)
            self.stats["total_documents"] += 1

            return ProcessedDocument(
                file_name=file_name,
                file_path=f"/sites/{notification.site_id}/lists/{notification.list_id}",
                site_name=notification.site_id,
                library_name=notification.list_id,
                mime_type=self._guess_mime(ext),
                size_bytes=len(content),
                text_content=text[:500],  # Preview only
                word_count=len(text.split()),
                drawer_ids=drawer_ids,
                status="processed",
            )

        except Exception as e:
            logger.error("Failed to process notification: %s", e)
            self.stats["total_errors"] += 1
            return ProcessedDocument(
                file_name=getattr(notification, "item_id", "unknown"),
                file_path="",
                site_name=notification.site_id,
                library_name=notification.list_id,
                mime_type="",
                size_bytes=0,
                text_content="",
                word_count=0,
                status="error",
                error=str(e),
            )

    async def process_drive(
        self,
        site_id: str,
        drive_id: str,
        folder_path: str = "/",
        site_name: str = "",
        drive_name: str = "",
    ) -> ProcessingBatch:
        """
        Process all documents in a SharePoint drive (full import).

        Args:
            site_id: SharePoint site ID
            drive_id: Document library (drive) ID
            folder_path: Path within the drive to process
            site_name: Human-readable site name (for MemPalace wing)
            drive_name: Human-readable library name (for MemPalace room)
        """
        batch = ProcessingBatch(
            site_id=site_id,
            site_name=site_name or site_id,
            drive_id=drive_id,
            drive_name=drive_name or drive_id,
            started_at=datetime.now(timezone.utc).isoformat(),
        )

        try:
            items = await self.graph.list_folder(site_id, drive_id, folder_path)
        except Exception as e:
            logger.error("Failed to list folder %s: %s", folder_path, e)
            batch.completed_at = datetime.now(timezone.utc).isoformat()
            return batch

        for item in items:
            # Key-presence check, not truthiness: Graph returns a folder
            # object that can be empty {} (falsy) but is still a folder.
            if "folder" in item:
                # Recurse into subfolders (full-import must cover the
                # whole library — fixed 2026-08-07: previously folders
                # were skipped, so sites with folder-organized content
                # (e.g. Mark 8: /Design, /Digital Assets Web) imported
                # ZERO documents).
                sub_path = f"{folder_path.rstrip('/')}/{item['name']}"
                await self._process_folder(
                    batch, site_id, drive_id, sub_path,
                    site_name, drive_name, depth=1,
                )
                continue

            await self._process_file(
                batch, site_id, drive_id, folder_path, item,
                site_name, drive_name,
            )

        batch.completed_at = datetime.now(timezone.utc).isoformat()
        return batch

    async def _process_folder(
        self,
        batch: ProcessingBatch,
        site_id: str,
        drive_id: str,
        folder_path: str,
        site_name: str,
        drive_name: str,
        depth: int = 0,
    ) -> None:
        """Recursively process a drive folder (bounded by MAX_FOLDER_DEPTH)."""
        if depth > MAX_FOLDER_DEPTH:
            logger.warning("Folder depth limit (%d) reached at %s", MAX_FOLDER_DEPTH, folder_path)
            return
        try:
            items = await self.graph.list_folder(site_id, drive_id, folder_path)
        except Exception as e:
            logger.error("Failed to list folder %s: %s", folder_path, e)
            return
        for item in items:
            if "folder" in item:
                sub_path = f"{folder_path.rstrip('/')}/{item['name']}"
                await self._process_folder(
                    batch, site_id, drive_id, sub_path,
                    site_name, drive_name, depth=depth + 1,
                )
                continue
            await self._process_file(
                batch, site_id, drive_id, folder_path, item,
                site_name, drive_name,
            )

    async def _process_file(
        self,
        batch: ProcessingBatch,
        site_id: str,
        drive_id: str,
        folder_path: str,
        item: dict,
        site_name: str,
        drive_name: str,
    ) -> None:
        """Download, extract, and mine a single drive file."""
        file_name = item.get("name", "unknown")
        ext = Path(file_name).suffix.lower()

        if ext not in self.SUPPORTED_EXTENSIONS:
            self.stats["total_skipped"] += 1
            return

        if item.get("size", 0) > self.max_file_size:
            self.stats["total_skipped"] += 1
            return

        try:
            content = await self.graph.download_file(
                site_id, drive_id, item["id"]
            )

            content_hash = hashlib.sha256(content).hexdigest()
            if content_hash in self._processed_hashes:
                self.stats["total_skipped"] += 1
                return
            self._processed_hashes.add(content_hash)

            text = self._extract_text(content, ext, file_name)
            if not text.strip():
                self.stats["total_skipped"] += 1
                return

            drawer_ids = await self._mine_to_mempalace(
                text=text,
                wing=self._sanitize_wing(site_name or site_id),
                room=self._sanitize_room(drive_name or drive_id),
                source_file=file_name,
            )

            doc = ProcessedDocument(
                file_name=file_name,
                file_path=self._join_path(folder_path, file_name),
                site_name=site_name,
                library_name=drive_name,
                mime_type=self._guess_mime(ext),
                size_bytes=len(content),
                text_content=text,
                word_count=len(text.split()),
                drawer_ids=drawer_ids,
            )
            batch.documents.append(doc)
            batch.total_processed += 1
            self.stats["total_processed"] += 1
            self.stats["total_documents"] += 1

        except Exception as e:
            logger.error("Failed to process %s: %s", file_name, e)
            batch.documents.append(ProcessedDocument(
                file_name=file_name,
                file_path=self._join_path(folder_path, file_name),
                site_name=site_name,
                library_name=drive_name,
                mime_type="",
                size_bytes=0,
                text_content="",
                word_count=0,
                status="error",
                error=str(e),
            ))
            batch.total_errors += 1
            self.stats["total_errors"] += 1

        batch.total_skipped = self.stats["total_skipped"]
        batch.completed_at = datetime.now(timezone.utc).isoformat()
        return batch

    # ---- Text Extraction ----

    def _extract_text(self, content: bytes, ext: str, file_name: str) -> str:
        """Extract text from document content based on file type."""
        ext = ext.lower()

        # Plain text formats
        if ext in {".txt", ".csv", ".md", ".yaml", ".yml", ".json",
                   ".py", ".js", ".ts", ".html", ".css", ".xml",
                   ".log", ".cfg", ".ini", ".toml"}:
            return self._extract_text_file(content)

        # Office formats
        if ext == ".docx" and DOCX_AVAILABLE:
            return self._extract_docx(content)

        if ext == ".pdf" and PYMUPDF_AVAILABLE:
            return self._extract_pdf(content)

        if ext == ".xlsx" and OPENPYXL_AVAILABLE:
            return self._extract_xlsx(content)

        if ext == ".pptx" and PPTX_AVAILABLE:
            return self._extract_pptx(content)

        # OpenDocument formats (LibreOffice native: odt/ods/odp) —
        # stdlib only, always available
        if ext in (".odt", ".ods", ".odp"):
            return self._extract_odf(content, kind=ext.lstrip("."))

        # Visio diagrams — stdlib only (ZIP + XML)
        if ext == ".vsdx":
            return self._extract_vsdx(content)

        # Video/audio — on-prem transcription (ffmpeg + whisper)
        if ext in (".mp4", ".mkv", ".mov", ".webm"):
            return self._extract_video(content)
        if ext in (".mp3", ".wav", ".m4a"):
            from threadweave.connectors.sharepoint.video import (
                transcribe_audio,
            )
            return transcribe_audio(content)

        return ""

    def _extract_text_file(self, content: bytes) -> str:
        """Extract text from plain text files."""
        try:
            return content.decode("utf-8")
        except UnicodeDecodeError:
            try:
                return content.decode("latin-1")
            except UnicodeDecodeError:
                return content.decode("utf-8", errors="replace")

    def _extract_docx(self, content: bytes) -> str:
        """Extract text from .docx files."""
        with tempfile.NamedTemporaryFile(
            suffix=".docx", dir=self.temp_dir, delete=False
        ) as tmp:
            tmp.write(content)
            tmp_path = tmp.name

        try:
            doc = docx.Document(tmp_path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n\n".join(paragraphs)
        except Exception as e:
            logger.error("Failed to extract docx: %s", e)
            return ""
        finally:
            try:
                os.unlink(tmp_path)
            except OSError:
                pass

    def _extract_pdf(self, content: bytes) -> str:
        """Extract text from PDF files."""
        import fitz

        doc = fitz.open(stream=content, filetype="pdf")
        try:
            pages = []
            for page in doc:
                pages.append(page.get_text())
            return "\n\n".join(pages)
        finally:
            doc.close()

    def _extract_odf(self, content: bytes, kind: str) -> str:
        """Extract text from OpenDocument files (odt/ods/odp).

        LibreOffice's native formats are ZIP containers with a
        content.xml inside (ISO 26300). Extracted with stdlib only:
        paragraphs (text:p), headings (text:h), and table cells
        (table:table-cell) are joined so the detector sees readable
        content. Handles LibreOffice's default save formats so Linux
        users are first-class.
        """
        import zipfile
        from xml.etree import ElementTree as ET

        try:
            with zipfile.ZipFile(io.BytesIO(content)) as zf:
                if "content.xml" not in zf.namelist():
                    return ""
                xml_data = zf.read("content.xml")
        except zipfile.BadZipFile:
            return ""

        try:
            root = ET.fromstring(xml_data)
        except ET.ParseError:
            return ""

        parts: list[str] = []
        # Paragraphs and headings, in document order
        for elem in root.iter():
            tag = elem.tag.split("}")[-1]
            if tag in ("p", "h"):
                text = "".join(elem.itertext()).strip()
                if text:
                    parts.append(text)
        # Table cells (ods): include cell values so spreadsheet
        # knowledge isn't lost (dedup prevents repeats)
        if kind == "ods":
            for elem in root.iter():
                if elem.tag.split("}")[-1] == "table-cell":
                    text = "".join(elem.itertext()).strip()
                    if text:
                        parts.append(text)
        return "\n\n".join(parts)

    def _extract_vsdx(self, content: bytes) -> str:
        """Extract text from Visio .vsdx files (shape text per page).

        .vsdx is a ZIP container (OOXML family): each page lives in
        visio/pages/pageN.xml, and shape labels live in <Shape><Text>.
        Only the <Text> content is taken — geometry and style
        properties are noise for the detector.
        """
        import zipfile
        from xml.etree import ElementTree as ET

        try:
            with zipfile.ZipFile(io.BytesIO(content)) as zf:
                page_names = sorted(
                    n for n in zf.namelist()
                    if n.startswith("visio/pages/page") and n.endswith(".xml")
                )
                if not page_names:
                    return ""
                pages = [zf.read(n) for n in page_names]
        except zipfile.BadZipFile:
            return ""

        parts: list[str] = []
        for idx, xml_data in enumerate(pages):
            try:
                root = ET.fromstring(xml_data)
            except ET.ParseError:
                continue
            page_parts: list[str] = []
            for elem in root.iter():
                # Text nodes anywhere in the shape tree (incl. nested)
                if elem.tag.split("}")[-1] == "Text":
                    text = "".join(elem.itertext()).strip()
                    if text:
                        page_parts.append(text)
            if page_parts:
                parts.append(f"[Page {idx + 1}]")
                parts.extend(page_parts)
        return "\n".join(parts)

    def _extract_video(self, content: bytes) -> str:
        """Transcribe a video file on-prem (ffmpeg + faster-whisper).

        Returns "" when transcription isn't available (no whisper
        package, no ffmpeg, or a bad file) — the file is then skipped
        like any other unsupported type. The transcript text goes
        through the normal detection/ingest path.
        """
        from threadweave.connectors.sharepoint.video import transcribe_video

        return transcribe_video(content)

    def _extract_xlsx(self, content: bytes) -> str:
        """Extract text from .xlsx workbooks (cell values per sheet).

        Preserves sheet names as section headers so the detector sees
        topic context, not a flat cell dump.
        """
        import io

        wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True,
                                    data_only=True)
        parts = []
        for ws in wb.worksheets:
            rows = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c).strip() for c in row if c is not None
                         and str(c).strip()]
                if cells:
                    rows.append(" | ".join(cells))
            if rows:
                parts.append(f"[Sheet: {ws.title}]\n" + "\n".join(rows))
        wb.close()
        return "\n\n".join(parts)

    def _extract_pptx(self, content: bytes) -> str:
        """Extract text from .pptx decks (slide shapes + speaker notes)."""
        import io

        prs = pptx.Presentation(io.BytesIO(content))
        parts = []
        for i, slide in enumerate(prs.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    texts.append(f"[notes] {notes}")
            if texts:
                parts.append(f"[Slide {i}]\n" + "\n".join(texts))
        return "\n\n".join(parts)

    # ---- MemPalace Integration ----

    # ---- ThreadWeave API Integration (central ingestion) ----

    async def _mine_to_mempalace(
        self,
        text: str,
        wing: str,
        room: str,
        source_file: str = "",
        author_id: str = "",
    ) -> list[str]:
        """Submit extracted text to the central ingestion pipeline."""
        import httpx

        try:
            async with httpx.AsyncClient(timeout=30.0) as client:
                resp = await client.post(
                    "http://localhost:8000/api/v1/ingest",
                    json={
                        "content": text,
                        "source": "sharepoint",
                        "tenant_id": "default",
                        "metadata": {
                            "wing": wing,
                            "room": room,
                            "title": source_file,
                            "source_file": source_file,
                            "author_id": author_id,
                        },
                    },
                )
                resp.raise_for_status()
                data = resp.json()
                return [data.get("id", "")] if data.get("should_save") else []

        except Exception as e:
            logger.error("Ingest API call failed: %s", e)
            return []

    # ---- Helpers ----

    @staticmethod
    def _chunk_text(text: str, max_chars: int = 4000) -> list[str]:
        """Split text into chunks at sentence boundaries."""
        if len(text) <= max_chars:
            return [text] if text.strip() else []

        chunks = []
        sentences = text.replace("\\n", " ").split(". ")
        current = ""

        for sentence in sentences:
            candidate = f"{current}. {sentence}" if current else sentence
            if len(candidate) > max_chars and current:
                chunks.append(current.strip())
                current = sentence
            else:
                current = candidate

        if current.strip():
            chunks.append(current.strip())

        return chunks

    @staticmethod
    def _join_path(folder_path: str, file_name: str) -> str:
        """Join a drive folder path with a file name, avoiding double slashes."""
        base = folder_path.rstrip("/")
        if not base:
            return f"/{file_name}"
        return f"{base}/{file_name}"

    @staticmethod
    def _sanitize_wing(name: str) -> str:
        """Convert SharePoint name to MemPalace wing name."""
        return name.lower().replace(" ", "_").replace("-", "_")[:64]

    @staticmethod
    def _sanitize_room(name: str) -> str:
        """Convert library name to MemPalace room name."""
        return name.lower().replace(" ", "_").replace("-", "_")[:64]

    @staticmethod
    def _extract_filename(notification: ChangeNotification) -> str:
        """Best-effort filename extraction from notification."""
        # The notification's resource path contains the item ID
        # We'd normally need a separate Graph call to get metadata
        return f"item_{notification.item_id}"

    @staticmethod
    def _guess_mime(ext: str) -> str:
        """Guess MIME type from extension."""
        mimes = {
            ".pdf": "application/pdf",
            ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            ".txt": "text/plain",
            ".csv": "text/csv",
            ".json": "application/json",
            ".md": "text/markdown",
            ".html": "text/html",
            ".py": "text/x-python",
            ".js": "text/javascript",
        }
        return mimes.get(ext, "application/octet-stream")
